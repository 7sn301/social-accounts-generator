from io import BytesIO
from datetime import datetime
from typing import Any, Dict, Iterable, List

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


def _safe(v: Any) -> str:
    if v is None:
        return ""
    if isinstance(v, (list, tuple, set)):
        return ", ".join(_safe(x) for x in v if x not in (None, ""))
    if isinstance(v, dict):
        return ", ".join(f"{k}: {_safe(val)}" for k, val in v.items() if val not in (None, ""))
    return str(v)


def _normalize_records(records: Iterable[Dict[str, Any]]) -> List[Dict[str, Any]]:
    return [r for r in (records or []) if isinstance(r, dict)]


def _rtl_paragraph(paragraph, text: str, bold: bool = False, size: int = 11):
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = "Arial"
    return paragraph


def _add_summary_table(doc: Document, records: List[Dict[str, Any]], columns: List[str], labels: Dict[str, str]):
    table = doc.add_table(rows=1, cols=len(columns))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = table.rows[0].cells
    for i, col in enumerate(columns):
        p = hdr[i].paragraphs[0]
        _rtl_paragraph(p, labels.get(col, col), bold=True, size=10)
    for rec in records:
        row = table.add_row().cells
        for i, col in enumerate(columns):
            p = row[i].paragraphs[0]
            _rtl_paragraph(p, _safe(rec.get(col, "")), size=9)


def generate_word_report(records: Iterable[Dict[str, Any]], title: str = "تقرير OSINT") -> bytes:
    records = _normalize_records(records)
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.5)
    section.bottom_margin = Inches(0.5)
    section.left_margin = Inches(0.55)
    section.right_margin = Inches(0.55)

    p = doc.add_paragraph()
    _rtl_paragraph(p, title, bold=True, size=18)
    p = doc.add_paragraph()
    _rtl_paragraph(p, f"تاريخ الإنشاء: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}", size=10)
    p = doc.add_paragraph()
    _rtl_paragraph(p, f"عدد السجلات: {len(records)}", size=10)

    if not records:
        p = doc.add_paragraph()
        _rtl_paragraph(p, "لا توجد بيانات لإدراجها في التقرير.", size=11)
    else:
        keys_priority = [
            "username", "user_screen_name", "user_name", "region_name_ar", "region_flag",
            "region_confidence", "region_source", "location", "country", "status", "tweet_url", "profile_url"
        ]
        seen = []
        for k in keys_priority:
            if any(k in r for r in records):
                seen.append(k)
        for k in records[0].keys():
            if k not in seen:
                seen.append(k)
        summary_cols = seen[: min(8, len(seen))]
        labels = {
            "username": "اسم المستخدم",
            "user_screen_name": "اسم المستخدم",
            "user_name": "الاسم",
            "region_name_ar": "الدولة",
            "region_flag": "العلم",
            "region_confidence": "الثقة",
            "region_source": "المصدر",
            "location": "الموقع",
            "country": "الدولة",
            "status": "الحالة",
            "tweet_url": "الرابط",
            "profile_url": "رابط الحساب",
        }
        p = doc.add_paragraph()
        _rtl_paragraph(p, "ملخص جدولي", bold=True, size=13)
        _add_summary_table(doc, records[:50], summary_cols, labels)

        doc.add_section(WD_SECTION.NEW_PAGE)
        p = doc.add_paragraph()
        _rtl_paragraph(p, "التفاصيل", bold=True, size=14)
        for i, rec in enumerate(records, 1):
            p = doc.add_paragraph()
            _rtl_paragraph(p, f"سجل رقم {i}", bold=True, size=12)
            for k, v in rec.items():
                if v in (None, "", [], {}, ()): 
                    continue
                p = doc.add_paragraph()
                _rtl_paragraph(p, f"{k}: {_safe(v)}", size=10)

    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _ppt_set_text(frame, text: str, level: int = 0, size: int = 20, bold: bool = False):
    p = frame.add_paragraph() if frame.text else frame.paragraphs[0]
    p.text = text
    p.level = level
    p.alignment = PP_ALIGN.RIGHT
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = PptPt(size)
        run.font.bold = bold
    return p


def generate_pptx_report(records: Iterable[Dict[str, Any]], title: str = "تقرير OSINT") -> bytes:
    records = _normalize_records(records)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"تاريخ الإنشاء: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}\nعدد السجلات: {len(records)}"

    def add_bullet_slide(slide_title: str, bullets: List[str]):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide_title
        tf = s.placeholders[1].text_frame
        tf.clear()
        first = True
        for bullet in bullets:
            if first:
                p = tf.paragraphs[0]
                p.text = bullet
                p.alignment = PP_ALIGN.RIGHT
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = PptPt(18)
                first = False
            else:
                _ppt_set_text(tf, bullet, size=18)

    if records:
        sample = records[: min(10, len(records))]
        overview = []
        if any(r.get("region_name_ar") for r in sample):
            countries = [str(r.get("region_name_ar")) for r in sample if r.get("region_name_ar")]
            if countries:
                top = {}
                for c in countries:
                    top[c] = top.get(c, 0) + 1
                ranked = sorted(top.items(), key=lambda x: x[1], reverse=True)[:5]
                overview.append("أكثر الدول ظهورًا: " + "، ".join(f"{k} ({v})" for k, v in ranked))
        high_conf = sum(1 for r in records if str(r.get("region_confidence", "0")).isdigit() and int(r.get("region_confidence", 0)) >= 70)
        overview.append(f"عدد النتائج عالية الثقة: {high_conf}")
        overview.append(f"عدد السجلات في التقرير: {len(records)}")
        add_bullet_slide("الملخص التنفيذي", overview)

        for idx, rec in enumerate(records[:20], 1):
            bullets = []
            for k, v in rec.items():
                if v in (None, "", [], {}, ()): 
                    continue
                bullets.append(f"{k}: {_safe(v)}")
                if len(bullets) >= 10:
                    break
            add_bullet_slide(f"سجل {idx}", bullets or ["لا توجد تفاصيل متاحة"])
    else:
        add_bullet_slide("الملخص التنفيذي", ["لا توجد بيانات لإدراجها في العرض"])

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()
